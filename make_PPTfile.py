def df_to_powerpoint(filename, df, **kwargs):
    """Converts a Pandas DataFrame to a table in a new, blank PowerPoint
    presentation.

    Creates a new PowerPoint presentation with the given filename, with a single
    slide containing a single table with the Pandas DataFrame data in it.

    The table is a standard Powerpoint table, and can easily be modified with
    the Powerpoint tools, for example: resizing columns, changing formatting
    etc.

    Parameters
    ----------
    filename: Filename to save the PowerPoint presentation as

    df: pandas ``DataFrame``
        DataFrame with the data

    **kwargs
        All other arguments that can be taken by ``df_to_table()`` (such as
        ``col_formatters`` or ``rounding``) can also be passed here.

    Returns
    -------
    pptx.shapes.graphfrm.GraphicFrame
        The python-pptx table (GraphicFrame) object that was created (which can
        then be used to do further manipulation if desired)
    """
    pres = Presentation()
    blank_slide_layout = pres.slide_layouts[6]
    slide = pres.slides.add_slide(blank_slide_layout)
    table = df_to_table(slide, df, **kwargs)
    pres.save(filename)

    return table 
  
  def ppt_section_slide(title, subtitle, file_path):
    from pptx import Presentation
    try:
        prs = Presentation(file_path)
    except:
        prs = Presentation(r'C:\Users\adrose\Desktop\AMD PowerPoint Template.pptx')
    picture_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(picture_slide_layout)

    # for x in slide.placeholders:
    #     print('%d %s' % (x.placeholder_format.idx, x.name))

    title_placeholder = slide.placeholders[0]
    subtitle_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = subtitle

    prs.save(file_path) 
