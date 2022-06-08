def pptx_to_text(document_path, event, context):
    import pptx

    prs = pptx.Presentation(document_path)

    paragraphs = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            paragraphs += [' '.join([run.text for run in paragraph.runs]) for paragraph in shape.text_frame.paragraphs]
        #end for
    #end for

    return '\n\n'.join(paragraphs).strip()
#end def 

def __init__(self, path_to_presentation, slide_layout=None, shape_properties=None):

        self.presentation  = Presentation(path_to_presentation) # TODO PptxPainter - Path checking # type: Presentation
        if slide_layout is None:
            self.default_slide_layout = None
        else:
            self.default_slide_layout = self.set_slide_layout(slide_layout)

        # Add all the dafault dicts to the class -
        if shape_properties:
            self._shape_properties = shape_properties
        else:
            self._shape_properties = PptxDefaults()

        self.textbox = self._shape_properties.textbox
        self.textbox_header = self._shape_properties.textbox_header
        self.textbox_footer = self._shape_properties.textbox_footer
        self.chart = self._shape_properties.chart
        self.table = self._shape_properties.table
        self.side_table = self._shape_properties.side_table

        charts = self._shape_properties.charts
        self.chart_bar = charts['bar']
        self.chart_bar_stacked100 = charts['bar_stacked100']
        self.chart_line = charts['line']
        self.chart_column = charts['column']
        self.chart_pie = charts['pie']

        self.slide_kwargs = {
                    'textboxs': {},
                    'charts': {},
                    'tables': {},
                    'side_tables': {},
                    } 
